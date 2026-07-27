"""Export targets for Sketch decks (kind: deck).

Sketch decks are authored in the Sketch tab as `kind: deck` analyses
with one Excalidraw scene per slot; PNG exports live under
`figures/<sketch-id>.png`. This module turns one of those decks into
a portable artefact:

  · to_pptx → vault/exports/<slug>.pptx (one image-fullbleed PPTX
    slide per slot, slot label as the slide title)
  · to_html → vault/exports/<slug>.html (single-file standalone
    reveal.js HTML with PNGs base64-embedded so the file is
    self-contained)

PDF export is deferred to Phase 5 (LibreOffice headless OR puppeteer
render — both add a heavy dependency we want to avoid until needed).
"""

from __future__ import annotations

import base64
import logging
import re
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from . import analyses

log = logging.getLogger("switchbay.deck_export")


def _exports_dir(workspace: Path) -> Path:
    d = workspace / "vault" / "exports"
    d.mkdir(parents=True, exist_ok=True)
    return d


def _label_for_slot(sketch_id: str) -> str:
    """Mirrors the SlidesTab label derivation so the same readable
    slot title is used wherever a deck is rendered."""
    parts = [p for p in re.split(r"[-_]+", sketch_id) if p]
    return " ".join(w[:1].upper() + w[1:] for w in parts) if parts else sketch_id


def _resolve_deck(workspace: Path, deck_path: str) -> dict:
    deck = analyses.load_analysis(workspace, deck_path)
    if deck is None:
        raise FileNotFoundError(f"deck not found: {deck_path}")
    if deck.get("kind") != "deck":
        raise ValueError(f"not a Sketch deck (kind={deck.get('kind')!r}): {deck_path}")
    return deck


def _png_for(workspace: Path, sketch_id: str) -> Path | None:
    # CE-native location first (2026-07-05 migration), legacy root after.
    for p in (
        workspace / "wiki" / "figures" / "_assets" / f"{sketch_id}.png",
        workspace / "figures" / f"{sketch_id}.png",
    ):
        if p.is_file():
            return p
    return None


def to_pptx(workspace: Path, deck_path: str) -> Path:
    """One PPTX slide per Sketch slot. Each slide is title + a
    centered, scaled-to-fit image of the slot's PNG. Missing PNGs
    fall through to a title-only slide rather than crashing the
    export — the user gets a deck they can open and the missing
    slot is obvious.
    """
    deck = _resolve_deck(workspace, deck_path)
    slug = str(deck["slug"])
    title = str(deck.get("title") or slug)

    prs = Presentation()
    # python-pptx defaults to 4:3 (10" x 7.5"). Bump to 16:9 since
    # Sketch slots are rendered against a 16:9 canvas in the Slides
    # tab and we want the export to match.
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    title_only_layout = prs.slide_layouts[5]

    slots = list(deck.get("slides") or [])
    notes = deck.get("slide_notes") or {}
    if not slots:
        slide = prs.slides.add_slide(title_only_layout)
        slide.shapes.title.text = title
        out = _exports_dir(workspace) / f"{slug}.pptx"
        prs.save(out)
        return out

    for sid in slots:
        sid_str = str(sid).strip()
        if not sid_str:
            continue
        slide = prs.slides.add_slide(title_only_layout)
        slide.shapes.title.text = _label_for_slot(sid_str)
        # Presenter note → the slide's notes page (PowerPoint speaker
        # view / "Notes" pane). Set even when the PNG is missing.
        note = str(notes.get(sid_str) or "").strip()
        if note:
            slide.notes_slide.notes_text_frame.text = note
        png = _png_for(workspace, sid_str)
        if png is None:
            continue
        # Centered image, leaving the title bar visible at the top.
        # PowerPoint's title placeholder occupies the upper ~1.25",
        # so anchor the image below that and let it fill the rest.
        slide.shapes.add_picture(
            str(png),
            Inches(0.4),
            Inches(1.4),
            width=Inches(12.5),
            height=Inches(5.7),
        )

    out = _exports_dir(workspace) / f"{slug}.pptx"
    prs.save(out)
    log.info("exported deck %s → %s", slug, out)
    return out


def to_html(workspace: Path, deck_path: str) -> Path:
    """Single-file reveal.js HTML with PNGs base64-embedded so the
    output works offline and round-trips by email/Slack. Uses the
    reveal CDN for the engine itself (keeps the file under ~1 MB
    even for 10-slot decks) — the local frontend bundles reveal.js
    too but we don't want to copy its CSS/JS into every export.
    """
    deck = _resolve_deck(workspace, deck_path)
    slug = str(deck["slug"])
    title = str(deck.get("title") or slug)
    slots = [str(s).strip() for s in (deck.get("slides") or []) if str(s).strip()]
    notes = deck.get("slide_notes") or {}

    def _note_aside(sid: str) -> str:
        # reveal's speaker view ('s' key) reads <aside class="notes">.
        note = str(notes.get(sid) or "").strip()
        return f"<aside class=\"notes\">{_h(note)}</aside>" if note else ""

    sections: list[str] = []
    if not slots:
        sections.append(f"<section><h1>{_h(title)}</h1></section>")
    else:
        for sid in slots:
            label = _label_for_slot(sid)
            png = _png_for(workspace, sid)
            if png is None:
                sections.append(
                    f"<section><h2>{_h(label)}</h2>"
                    f"<p><em>(missing figure)</em></p>{_note_aside(sid)}</section>"
                )
                continue
            b64 = base64.b64encode(png.read_bytes()).decode("ascii")
            sections.append(
                f"<section><h2>{_h(label)}</h2>"
                f"<img src=\"data:image/png;base64,{b64}\" "
                f"alt=\"{_h(label)}\" style=\"max-width:100%;max-height:80vh;\" />"
                f"{_note_aside(sid)}</section>"
            )

    html = _HTML_TEMPLATE.format(
        title=_h(title),
        sections="\n".join(sections),
    )
    out = _exports_dir(workspace) / f"{slug}.html"
    out.write_text(html, encoding="utf-8")
    log.info("exported deck %s → %s", slug, out)
    return out


def _h(s: str) -> str:
    return (
        s.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace("\"", "&quot;")
    )


_HTML_TEMPLATE = """<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8" />
  <title>{title}</title>
  <link rel="stylesheet"
    href="https://cdn.jsdelivr.net/npm/reveal.js@5.2.1/dist/reset.css" />
  <link rel="stylesheet"
    href="https://cdn.jsdelivr.net/npm/reveal.js@5.2.1/dist/reveal.css" />
  <link rel="stylesheet"
    href="https://cdn.jsdelivr.net/npm/reveal.js@5.2.1/dist/theme/black.css" />
</head>
<body>
  <div class="reveal"><div class="slides">
    {sections}
  </div></div>
  <script src="https://cdn.jsdelivr.net/npm/reveal.js@5.2.1/dist/reveal.js"></script>
  <script src="https://cdn.jsdelivr.net/npm/reveal.js@5.2.1/dist/plugin/notes/notes.js"></script>
  <script>Reveal.initialize({{ hash: true, controls: true, progress: true, plugins: [ RevealNotes ] }});</script>
</body>
</html>
"""
